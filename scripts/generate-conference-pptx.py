#!/usr/bin/env python3
"""Build a polished 25-slide international conference PowerPoint for PrognosEd."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "docs" / "print-assets"
OUT = ROOT / "docs" / "PrognosEd_International_Conference_Presentation.pptx"

PRIMARY = RGBColor(0x0B, 0x6E, 0x4F)
CHARCOAL = RGBColor(0x10, 0x18, 0x28)
MUTED = RGBColor(0x47, 0x54, 0x67)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF7, 0xF5)
SOFT = RGBColor(0xE6, 0xF4, 0xEE)
TOTAL = 25


def set_run(run, size=18, bold=False, color=CHARCOAL, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def send_back(slide, shape):
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    send_back(slide, shape)


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_top_rule(slide):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.03)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = SOFT
    rule.line.fill.background()


def add_footer(slide, page):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(9.5), Inches(0.3))
    p = left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"PrognosEd  ·  International Conference  ·  {page}/{TOTAL}"
    set_run(r, 11, False, MUTED)

    right = slide.shapes.add_textbox(Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.3))
    p = right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "AI · Student Success"
    set_run(r, 11, False, MUTED)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def title_box(slide, text, size=28):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size, True, PRIMARY)
    add_top_rule(slide)
    return box


def bullets(slide, items, left=0.55, top=1.25, width=12.2, height=5.4, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            item = {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item.get("level", 0)
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = item["text"]
        set_run(r, item.get("size", size), item.get("bold", False), item.get("color", CHARCOAL))
    return box


def content_slide(prs, heading, page, items, notes="", size=17):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    title_box(slide, heading)
    bullets(slide, items, size=size)
    add_footer(slide, page)
    if notes:
        add_notes(slide, notes)
    return slide


def picture_slide(prs, heading, image_path, caption, page, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    title_box(slide, heading, size=26)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1.35), Inches(1.2), width=Inches(10.6))
    cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    set_run(r, 12, False, MUTED)
    add_footer(slide, page)
    if notes:
        add_notes(slide, notes)
    return slide


def metric_card(slide, left, top, title, value, subtitle):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.8), Inches(1.55)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = SOFT

    t = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(5.3), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run(r, 13, True, PRIMARY)

    v = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.52), Inches(5.3), Inches(0.5))
    r = v.text_frame.paragraphs[0].add_run()
    r.text = value
    set_run(r, 26, True, CHARCOAL)

    s = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 1.05), Inches(5.3), Inches(0.35))
    r = s.text_frame.paragraphs[0].add_run()
    r.text = subtitle
    set_run(r, 12, False, MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- 1 Title ----
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()
    bottom = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = PRIMARY
    bottom.line.fill.background()

    box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "PrognosEd"
    set_run(r, 48, True, PRIMARY)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    r = p.add_run()
    r.text = "AI-Powered Early-Warning Student Success Platform"
    set_run(r, 26, True, CHARCOAL)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(16)
    r = p.add_run()
    r.text = "Predictive Analytics  ·  Intervention Intelligence  ·  Retrieval-Augmented Advising"
    set_run(r, 15, False, MUTED)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(28)
    r = p.add_run()
    r.text = "International Conference Presentation"
    set_run(r, 16, True, CHARCOAL)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    r = p.add_run()
    r.text = "[Author Name]  ·  [Affiliation]  ·  [Conference / Year]"
    set_run(r, 14, False, MUTED)
    add_notes(
        s,
        "Welcome the audience. Introduce PrognosEd as an AI platform for early academic risk detection, "
        "personalized interventions, and faculty decision support.",
    )

    # ---- 2 Outline ----
    content_slide(
        prs,
        "Presentation Outline",
        2,
        [
            "1. Problem context and research motivation",
            "2. Objectives, proposed solution, and system architecture",
            "3. AI/ML methodology: four specialized predictive models",
            "4. Key innovation — retrieval-augmented (RAG) advising chat",
            "5. Role-based decision support and live prototype demonstration",
            "6. Experimental results, impact, limitations, and future work",
            "7. Conclusions and discussion",
        ],
        "Keep this slide short. It orients the audience for a 15–20 minute talk.",
    )

    # ---- 3 Challenge ----
    content_slide(
        prs,
        "The Challenge: Delayed Intervention in Higher Education",
        3,
        [
            "Universities often detect academic distress after grades have already declined.",
            "Faculty workload limits continuous monitoring across large course cohorts.",
            "Attendance, LMS activity, and assessment signals remain fragmented and underused.",
            "One-size-fits-all advising fails to prioritize who needs support first.",
            "Result: reactive remediation, higher dropout risk, and weaker retention outcomes.",
            "Need: an early-warning intelligence system that predicts risk and recommends action.",
        ],
        "Stress timing: earlier signals enable lower-cost, higher-impact support.",
    )

    # ---- 4 Objectives ----
    content_slide(
        prs,
        "Research Objectives and Scope",
        4,
        [
            "Design an AI-enabled platform for early academic and dropout risk prediction.",
            "Deliver personalized, ranked intervention recommendations for faculty action.",
            "Provide role-based dashboards for Faculty, Admin, Department Head, and Director.",
            "Integrate multi-model ML services with graceful offline degradation.",
            "Introduce conversational RAG advising grounded in student metrics and history.",
            "Validate through a production-structured prototype with measurable model metrics.",
        ],
        "Frame each objective as an engineering and research outcome.",
    )

    # ---- 5 Solution ----
    content_slide(
        prs,
        "Proposed Solution Overview",
        5,
        [
            "PrognosEd — an integrated Student Success Intelligence Platform.",
            "Ingests student engagement and academic metrics (attendance, GPA, LMS, late work).",
            "Predicts risk via four independent ML microservices.",
            "Surfaces prioritized alerts and intervention plans in role-scoped dashboards.",
            "Supports decisions with a per-student RAG advising assistant.",
            "Secured with JWT authentication, RBAC, rate limiting, and SMTP faculty onboarding.",
        ],
        "Transition clearly from problem to one end-to-end system narrative.",
    )

    # ---- 6 Architecture (image preferred) ----
    img = ASSETS / "01-system-architecture.png"
    if img.exists():
        picture_slide(
            prs,
            "System Architecture",
            img,
            "Four-tier architecture: Presentation · Application · AI & Analytics · Knowledge",
            6,
            "Walk top-to-bottom: UI → API/DB → ML services → RAG knowledge layer.",
        )
    else:
        content_slide(
            prs,
            "System Architecture",
            6,
            [
                "Presentation Layer: React 19 SPA, role-based dashboards, RAG chat UI",
                "Application Layer: Express REST API, JWT sessions, SQLite persistence",
                "AI & Analytics Layer: FastAPI ML services on ports 8000–8003",
                "Knowledge Layer: TF-IDF retriever, advising playbook, chat memory",
                "Design principle: modular services, clear APIs, resilient fallbacks",
            ],
        )

    # ---- 7 Tech stack ----
    content_slide(
        prs,
        "Technology Stack",
        7,
        [
            "Frontend: React 19, Vite, Tailwind CSS 4, Recharts",
            "Backend: Node.js, Express 4, better-sqlite3, bcrypt, JWT",
            "ML Serving: FastAPI + uvicorn, scikit-learn / XGBoost pipelines",
            "Advising AI: Node.js TF-IDF RAG retriever with intent-aware generation",
            "Ops: Docker Compose support, GitHub Actions CI (lint, build, smoke tests)",
            "Email: Nodemailer SMTP for invitations and password reset",
        ],
        "Highlight production-structured choices beyond research notebooks.",
    )

    # ---- 8 Pipeline (image preferred) ----
    img = ASSETS / "02-ai-workflow.png"
    if img.exists():
        picture_slide(
            prs,
            "AI Pipeline: From Student Metrics to Action",
            img,
            "Metrics → Risk Prediction → Intervention Ranking → Faculty Action → RAG Advising",
            8,
            "Emphasize the closed loop from data to action to continuous advising.",
        )
    else:
        content_slide(
            prs,
            "AI Pipeline: From Student Metrics to Action",
            8,
            [
                "Step 1 — Collect metrics: attendance, GPA, LMS activity, late assignments",
                "Step 2 — Predict risk: academic early-warning + dropout probability",
                "Step 3 — Rank interventions: skill gaps + risk-level action templates",
                "Step 4 — Faculty action: dashboard alerts, accept/dismiss recommendations",
                "Step 5 — Continuous advising: RAG chat with per-student conversation memory",
                "Feedback loop: recommendation decisions refine operational prioritization",
            ],
        )

    # ---- 9 Model 1 ----
    content_slide(
        prs,
        "Model 1 — Academic Early-Warning Risk Prediction",
        9,
        [
            "Purpose: predict pass/fail risk early enough for meaningful intervention",
            "Signals: LMS engagement (clicks, active days) and assessment performance",
            "Evaluation design: 4-week early-warning window (not full-course leakage)",
            "Key result: ROC-AUC ≈ 0.83 on OULAD early-warning setting",
            "Explainability: SHAP confirms behavioral/academic drivers over demographics",
            "Service endpoint: FastAPI on port 8000 (ML_ACADEMIC_API_URL)",
        ],
        "Stress the 4-week constraint as methodological honesty for real-world timing.",
    )

    # ---- 10 Model 2 ----
    content_slide(
        prs,
        "Model 2 — Dropout Risk Detection",
        10,
        [
            "Purpose: binary classification of Dropout vs Graduate outcomes",
            "Features: enrollment, curricular units, financial/academic indicators",
            "Dataset: UCI Higher Education Students Performance / Dropout",
            "Key results: 94.1% accuracy; AUC 0.973",
            "Validation: stratified 5-fold CV + held-out test reporting",
            "Service endpoint: FastAPI on port 8001 (ML_DROPOUT_API_URL)",
        ],
        "Report both accuracy and AUC; mention CV to avoid single-split optimism.",
    )

    # ---- 11 Model 3 ----
    content_slide(
        prs,
        "Model 3 — Skill-Based Intervention Recommender",
        11,
        [
            "Purpose: recommend skill-level interventions from mastery gaps",
            "Approach: rank weak skill areas from tutoring / EDM mastery tables",
            "Dataset lineage: KDD Cup 2010 Educational Data Mining challenge",
            "Key result: Precision@5 = 100% on evaluated ranking setting",
            "Integration: complements platform rule-based intervention templates",
            "Service endpoint: FastAPI on port 8002 (ML_RECOMMENDER_API_URL)",
        ],
        "Clarify that Model 3 ranks actionable skill targets, not only risk scores.",
    )

    # ---- 12 Model 4 ----
    content_slide(
        prs,
        "Model 4 — Behavioral Learning Analytics Clustering",
        12,
        [
            "Purpose: unsupervised segmentation for institutional analytics",
            "Method: k-means behavioral clustering (k = 4)",
            "Key metric: Silhouette score ≈ 0.315 (moderate, exploratory separation)",
            "Use case: cohort patterns for directors and academic administrators",
            "Supports reporting beyond individual risk lists",
            "Service endpoint: FastAPI on port 8003 (ML_ANALYTICS_API_URL)",
        ],
        "Frame clustering as exploratory analytics for leadership, not diagnosis.",
    )

    # ---- 13 RAG ----
    content_slide(
        prs,
        "Key Innovation: Retrieval-Augmented Advising Chat",
        13,
        [
            "Per-student conversational assistant embedded in the student profile",
            "Retrieves from an advising playbook using TF-IDF (Node.js, zero heavy deps)",
            "Personalizes answers with live metrics: risk, attendance, GPA, LMS, late work",
            "Stores short conversation memory per student–advisor pair in SQLite",
            "Intent detection: outreach drafts, weekly plans, risk explanation, summaries",
            "Anti-repetition controls rotate focus and avoid recycled playbook chunks",
        ],
        "Position RAG as the bridge between model outputs and faculty practice.",
    )

    # ---- 14 Roles ----
    content_slide(
        prs,
        "Role-Based Decision Support (Faculty → Leadership)",
        14,
        [
            "Faculty: course KPIs, at-risk roster, recommendations, RAG advisor",
            "Department Head: department performance and faculty-aligned views",
            "Academic Admin: institutional insights, user invites, course assignments",
            "Director / Dean: university KPIs, department comparison, risk drill-down",
            "Administrative Staff: reports and alerts for operational support",
            "Access control enforced via JWT + role-scoped service queries",
        ],
        "Show that AI insights are usable across governance levels.",
    )

    # ---- 15 Dashboard demo ----
    picture_slide(
        prs,
        "Prototype Demonstration — Faculty Dashboard",
        ASSETS / "05-faculty-dashboard.png",
        "Live UI: KPI cards, weekly engagement trend, and students needing attention",
        15,
        "Narrate a faculty morning workflow: scan KPIs → open at-risk list.",
    )

    # ---- 16 Student + RAG demo ----
    picture_slide(
        prs,
        "Prototype Demonstration — Student Risk Profile & Interventions",
        ASSETS / "06-student-profile-rag.png",
        "Live UI: critical-risk profile, intervention plan, and personalized RAG advisor",
        16,
        "Walk one student story end-to-end (high late work + low GPA → actions).",
    )

    # ---- 17 Results (cards + optional image note) ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_accent_bar(s)
    title_box(s, "Experimental Results and Model Performance")
    metric_card(s, 0.7, 1.35, "Model 1 · Academic Early-Warning", "ROC-AUC ≈ 0.83", "4-week OULAD early-warning setting")
    metric_card(s, 6.85, 1.35, "Model 2 · Dropout Risk", "94.1% Acc  ·  AUC 0.973", "UCI Dropout · stratified 5-fold CV")
    metric_card(s, 0.7, 3.2, "Model 3 · Skill Recommender", "Precision@5 = 100%", "KDD Cup 2010 mastery-gap ranking")
    metric_card(s, 6.85, 3.2, "Model 4 · Behavioral Clustering", "Silhouette ≈ 0.315", "k-means · k = 4 exploratory segments")
    extra = s.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(12), Inches(1.4))
    tf = extra.text_frame
    tf.word_wrap = True
    for i, line in enumerate(
        [
            "• Platform resilience: heuristic fallbacks when ML services are offline",
            "• Prototype coverage: dashboards, alerts, recommendations, reports, and RAG chat",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        set_run(r, 16, False, CHARCOAL)
    add_footer(s, 17)
    add_notes(s, "Lead with Model 2 headline metrics, then explain early-warning honesty for Model 1.")

    # ---- 18 Compliance ----
    content_slide(
        prs,
        "Specification Compliance and Deliverables",
        18,
        [
            "✓ Faculty dashboard with at-risk prioritization",
            "✓ Risk alert system with role-scoped visibility",
            "✓ Recommendation engine (rules + ML skill diagnostics)",
            "✓ Institutional / class reporting analytics",
            "✓ Four ML microservices for prediction and analytics",
            "✓ Extended deliverables: RAG chat, SMTP onboarding, RBAC hierarchy",
        ],
        "Useful for evaluators mapping against the original project specification.",
    )

    # ---- 19 Security ----
    content_slide(
        prs,
        "Security, Scalability, and Deployment",
        19,
        [
            "Security: JWT sessions, bcrypt hashing, API rate limits, CORS controls",
            "Privacy posture: role-scoped data access; demo accounts disabled in production",
            "Deployment: npm production build, Express static hosting, Docker Compose",
            "Ops readiness: health endpoints for API, ML services, and email",
            "Scale path: SQLite for pilot → PostgreSQL for institutional concurrency",
            "CI quality gates: lint, build, and smoke tests on push/PR",
        ],
        "Reassure audiences that the prototype is deployment-aware.",
    )

    # ---- 20 Limitations ----
    content_slide(
        prs,
        "Limitations and Ethical Considerations",
        20,
        [
            "Some ML features are estimated from summary metrics (not raw LMS event streams).",
            "RAG uses a curated playbook; not a full institutional policy corpus or LLM.",
            "Clustering is exploratory; moderate silhouette implies interpretive caution.",
            "Predictions must support—not replace—professional academic judgment.",
            "Fairness vigilance: minimize reliance on sensitive demographic proxies.",
            "Production use requires FERPA/GDPR-aligned governance and audit logging.",
        ],
        "Honesty builds credibility with international academic audiences.",
    )

    # ---- 21 Impact ----
    content_slide(
        prs,
        "Impact and Contributions",
        21,
        [
            "Shifts advising from reactive remediation to proactive early warning",
            "Unifies prediction, recommendation, and conversational guidance in one system",
            "Provides leadership visibility from course level to institutional KPIs",
            "Demonstrates modular, production-structured ML microservice integration",
            "Introduces per-student RAG memory for continuity of advising context",
            "Offers a deployable prototype pathway for university pilot programs",
        ],
        "Close the value loop: technical novelty + institutional usefulness.",
    )

    # ---- 22 Future ----
    content_slide(
        prs,
        "Future Work and Research Directions",
        22,
        [
            "Live LMS/SIS ingestion (LTI, xAPI, or scheduled ETL pipelines)",
            "Embedding-based retrieval + grounded LLM synthesis with citations",
            "PostgreSQL migration, connection pooling, and horizontal scaling",
            "Automated retraining, model versioning, and monitoring (e.g., MLflow)",
            "Real-time alert digests and institutional workflow integrations",
            "Expanded fairness audits, accessibility (WCAG), and longitudinal impact studies",
        ],
        "Invite collaboration: data partnerships, pilots, and joint research.",
    )

    # ---- 23 Conclusion ----
    content_slide(
        prs,
        "Conclusion",
        23,
        [
            "PrognosEd operationalizes AI for early student-success intervention.",
            "Four ML services cover academic risk, dropout, skills, and behavioral analytics.",
            "RAG advising adds personalized, history-aware decision support for faculty.",
            "Role-based dashboards connect course action to institutional strategy.",
            "Strong prototype readiness with clear, evidence-based performance metrics.",
            "Next step: pilot deployment with live institutional data partnerships.",
        ],
        "Restate the contribution in one sentence before Q&A.",
    )

    # ---- 24 Acknowledgments ----
    content_slide(
        prs,
        "Acknowledgments",
        24,
        [
            "Academic supervision and research guidance for the AI/ML internship track",
            "Open datasets enabling reproducible model development",
            {"text": "OULAD · UCI Student Performance · UCI Dropout · KDD Cup 2010 EDM", "level": 1, "size": 15},
            "Open-source ecosystem: React, Express, FastAPI, scikit-learn, and related tools",
            "Institutional stakeholders who inspired practical faculty advising workflows",
        ],
        "Keep acknowledgments concise and respectful.",
    )

    # ---- 25 Q&A ----
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Questions & Discussion"
    set_run(r, 42, True, PRIMARY)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(18)
    r = p.add_run()
    r.text = "Thank you for your attention"
    set_run(r, 22, False, CHARCOAL)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = "PrognosEd — AI-Powered Early-Warning Student Success Platform"
    set_run(r, 14, False, MUTED)
    add_footer(s, 25)
    add_notes(
        s,
        "Prepare concise answers on datasets, fairness, RAG limits, and pilot readiness.",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    from pptx import Presentation as P

    n = len(P(str(path)).slides)
    print(f"Wrote {path}")
    print(f"Slides: {n}")
